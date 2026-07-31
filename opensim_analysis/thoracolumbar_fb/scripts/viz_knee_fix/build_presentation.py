"""SMA 허리 보조 슈트 5동작 연구 발표 deck (33 slides, 영상 5종 임베드).

python-pptx로 생성. 영상은 add_movie()로 바이너리 임베드 (경로 링크 아님).
모든 수치는 /data/*_results/ SO 산출물에서 재계산·검증된 값.
"""
import os, math
from pptx import Presentation
from pptx.util import Inches as I, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from PIL import Image

# ---------------------------------------------------------------- 상수
NAVY   = RGBColor(0x1B, 0x3A, 0x5C)
DARK   = RGBColor(0x33, 0x33, 0x33)
RED    = RGBColor(0xC0, 0x39, 0x2B)
GREEN  = RGBColor(0x1E, 0x84, 0x49)
GRAY   = RGBColor(0x88, 0x88, 0x88)
LGRAY  = RGBColor(0xBB, 0xBB, 0xBB)
BOXBG  = RGBColor(0xF4, 0xF6, 0xF8)
BOXBG2 = RGBColor(0xEC, 0xF3, 0xEE)
BOXBG3 = RGBColor(0xFA, 0xEE, 0xEC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
FONT   = 'Malgun Gothic'

SW, SH = 13.333, 7.5
M      = 0.6                 # 최소 여백
CX     = SW / 2

MEDIA  = '/data/opensim_results/ppt_media'
IMG    = '/data/wearable-assist/opensim_analysis/thoracolumbar_fb/docs/images'
OUT    = '/data/opensim_results/SMA_suit_5motion_presentation.pptx'

prs = Presentation()
prs.slide_width  = I(SW)
prs.slide_height = I(SH)
BLANK = prs.slide_layouts[6]

_page = {'n': 0}


# ---------------------------------------------------------------- 헬퍼
def tb(slide, l, t, w, h, wrap=True):
    """빈 텍스트박스. auto_size 끄고 고정 크기 → 오버플로 제어를 코드가 담당."""
    s = slide.shapes.add_textbox(I(l), I(t), I(w), I(h))
    f = s.text_frame
    f.word_wrap = wrap
    f.margin_left = f.margin_right = Pt(0)
    f.margin_top = f.margin_bottom = Pt(0)
    return s, f


def _p(f, first):
    """first=True여도 paragraphs[0]에 이미 run이 있으면 새 문단을 만든다.
    (v1 결함: 헤더 뒤 bullets 첫 항목이 같은 문단에 붙어 '라벨•항목'으로 병합됨)"""
    if first and not f.paragraphs[0].runs:
        return f.paragraphs[0]
    return f.add_paragraph()


def para(f, text, size=17, bold=False, color=DARK, align=PP_ALIGN.LEFT,
         space_after=6, space_before=0, first=False, line=1.25, font=FONT):
    p = _p(f, first)
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.line_spacing = line
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return p


def rich(f, parts, size=17, align=PP_ALIGN.LEFT, first=False,
         space_after=6, space_before=0, line=1.25):
    """parts = [(text, bold, color, size|None), ...] 한 문단 내 서식 혼합."""
    p = _p(f, first)
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.line_spacing = line
    for t, b, c, sz in parts:
        r = p.add_run()
        r.text = t
        r.font.size = Pt(sz or size)
        r.font.bold = b
        r.font.color.rgb = c
        r.font.name = FONT
    return p


def slide(title=None, sub=None):
    s = prs.slides.add_slide(BLANK)
    _page['n'] += 1
    if title:
        _, f = tb(s, M, 0.42, SW - 2 * M, 0.85)
        para(f, title, size=32, bold=True, color=NAVY, first=True,
             space_after=0, line=1.05)
    if sub:
        _, f2 = tb(s, M, 1.22, SW - 2 * M, 0.4)
        para(f2, sub, size=15, color=GRAY, first=True, space_after=0)
    # 페이지 번호
    _, fp = tb(s, SW - M - 0.9, SH - 0.55, 0.9, 0.3)
    para(fp, str(_page['n']), size=10, color=GRAY,
         align=PP_ALIGN.RIGHT, first=True, space_after=0)
    return s


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


_boxes = []   # (shape, width_in, height_in) — 마지막에 일괄 autofit


def _est_height(tf, w_in):
    """텍스트 프레임이 실제로 차지할 높이(inch) 추정.
    한글 1.0 em, ASCII 0.52 em 폭으로 줄 수를 세고 행간·문단간격을 더한다."""
    avail = w_in - 0.30
    if avail <= 0:
        return 0.0
    tot = 0.0
    for p in tf.paragraphs:
        if not p.runs:
            continue
        sz = max((r.font.size.pt if r.font.size else 17) for r in p.runs)
        em = sz / 72.0
        txt = ''.join(r.text for r in p.runs)
        wdt = sum(1.0 if ord(c) > 0x2E80 else 0.52 for c in txt) * em
        nl = max(1, math.ceil(wdt / avail - 1e-6))
        ls = p.line_spacing if isinstance(p.line_spacing, float) else 1.25
        tot += nl * em * 1.22 * ls
        tot += (p.space_after.pt if p.space_after else 0) / 72.0
        tot += (p.space_before.pt if p.space_before else 0) / 72.0
    return tot


def autofit_all(margin=0.20, floor=0.62):
    """등록된 모든 박스를 검사해 넘치면 폰트를 단계적으로 축소한다.
    (고정 높이 박스에서 텍스트가 테두리를 뚫고 나가는 결함의 근본 차단)"""
    shrunk = []
    for sp, w_in, h_in in _boxes:
        tf = sp.text_frame
        cap = h_in - margin
        if cap <= 0 or _est_height(tf, w_in) <= cap:
            continue
        sc = 1.0
        while sc > floor and _est_height(tf, w_in) > cap:
            sc -= 0.04
            for p in tf.paragraphs:
                for r in p.runs:
                    if r.font.size:
                        r.font.size = Pt(round(r.font.size.pt * 0.96, 1))
        shrunk.append((round(w_in, 2), round(h_in, 2), round(sc, 2)))
    return shrunk


def rect(slide, l, t, w, h, fill=BOXBG, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         adj=0.06):
    sp = slide.shapes.add_shape(shape, I(l), I(t), I(w), I(h))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = adj
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(10)
    tf.margin_top = tf.margin_bottom = Pt(8)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    _boxes.append((sp, w, h))
    return sp, tf


def arrow(slide, l, t, w, h, color=NAVY):
    sp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, I(l), I(t), I(w), I(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def down_arrow(slide, l, t, w, h, color=NAVY):
    sp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, I(l), I(t), I(w), I(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def fit(path, boxw, boxh):
    """종횡비 유지하며 (boxw, boxh) 안에 들어가는 (w, h) 반환."""
    with Image.open(path) as im:
        iw, ih = im.size
    sc = min(boxw / iw, boxh / ih)
    return iw * sc, ih * sc


def pic(slide, path, cl, ct, boxw, boxh, center=True):
    """boxw x boxh 영역 중앙에 종횡비 유지 배치."""
    w, h = fit(path, boxw, boxh)
    l = cl + (boxw - w) / 2 if center else cl
    t = ct + (boxh - h) / 2
    return slide.shapes.add_picture(path, I(l), I(t), I(w), I(h))


def movie(slide, key, boxw=9.0, boxh=4.45, top=1.72):
    """영상 임베드 (poster_frame_image 필수 — 없으면 검은 사각형)."""
    v = f'{MEDIA}/{key}_ppt.mp4'
    p = f'{MEDIA}/poster_{key}.png'
    w, h = fit(p, boxw, boxh)
    l = CX - w / 2
    return slide.shapes.add_movie(v, I(l), I(top), I(w), I(h),
                                  poster_frame_image=p, mime_type='video/mp4')


def headline(slide, parts, top=1.60, size=40, align=PP_ALIGN.CENTER, h=0.85):
    _, f = tb(slide, M, top, SW - 2 * M, h)
    rich(f, parts, size=size, align=align, first=True, space_after=0, line=1.05)


def caption(slide, text, top=6.30, size=14, color=GRAY, align=PP_ALIGN.CENTER,
            h=0.6, bold=False):
    """페이지 번호(우하단)와 겹치지 않도록 폭을 좁혀 중앙 배치."""
    w = SW - 2 * M - 1.3
    _, f = tb(slide, CX - w / 2, top, w, h)
    para(f, text, size=size, color=color, align=align, first=True,
         space_after=0, bold=bold, line=1.3)


def bullets(f, items, size=17, color=DARK, sa=9, indent_marker='• '):
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            txt, bold, col = it
        else:
            txt, bold, col = it, False, color
        para(f, indent_marker + txt, size=size, bold=bold, color=col,
             first=(i == 0), space_after=sa)


VLEGEND = ('왼쪽 = 슈트 없음  |  오른쪽 = 슈트 착용      '
           '근육 색:  초록(편함) → 노랑 → 빨강(힘듦)')


# ================================================================ S1 표지
s = prs.slides.add_slide(BLANK)
_page['n'] += 1
_, f = tb(s, 0.6, 2.05, SW - 1.2, 1.9)
para(f, 'SMA 직물 근육 액추에이터 기반 허리 보조 웨어러블 슈트의',
     size=28, bold=True, color=NAVY, align=PP_ALIGN.CENTER, first=True,
     space_after=6, line=1.18)
para(f, '근골격계 시뮬레이션 기반 효과 검증',
     size=28, bold=True, color=NAVY, align=PP_ALIGN.CENTER, space_after=0, line=1.18)
_, f = tb(s, 1.0, 4.12, SW - 2.0, 0.6)
para(f, '5가지 작업 동작에서의 척추기립근 부하 저감 정량 평가',
     size=20, color=DARK, align=PP_ALIGN.CENTER, first=True, space_after=0)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(5.17), I(4.95), I(3.0), Emu(9525 * 2))
ln.fill.solid(); ln.fill.fore_color.rgb = LGRAY; ln.line.fill.background()
ln.shadow.inherit = False
_, f = tb(s, 1.0, 5.35, SW - 2.0, 1.0)
para(f, '한국기계연구원 (KIMM)', size=16, color=DARK,
     align=PP_ALIGN.CENTER, first=True, space_after=4)
para(f, '발표자: 박철훈    |    2026. 07.', size=16, color=GRAY,
     align=PP_ALIGN.CENTER, space_after=0)
notes(s, '본 연구는 SMA 직물 근육을 사용한 허리 보조 슈트가 실제 작업 동작에서 '
         '허리 근육 부담을 얼마나 줄이는지를, 전신 근골격계 시뮬레이션으로 '
         '정량 평가하고 시각화한 결과입니다.')

# ================================================================ S2 목차
s = slide('목차')
items = ['연구 배경 및 목표', 'SMA 허리 보조 슈트 개요', '시뮬레이션 방법론',
         '5가지 작업 동작별 결과', '통합 결과 — 부하·효과 패턴', '고찰',
         '결론 및 향후 계획']
top = 1.75
for i, it in enumerate(items):
    num, nf = rect(s, 2.6, top, 0.62, 0.62, fill=NAVY)
    para(nf, str(i + 1), size=20, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, first=True, space_after=0)
    nf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _, tf = tb(s, 3.45, top + 0.11, 7.0, 0.5)
    para(tf, it, size=21, color=DARK, first=True, space_after=0)
    top += 0.75
notes(s, '오늘 발표는 크게 일곱 부분으로 구성됩니다. 배경과 슈트 개요를 짧게 보고, '
         '방법론을 설명한 뒤 다섯 동작 결과를 하나씩 보여드리겠습니다. '
         '핵심은 다섯 번째, 통합 결과입니다.')

# ================================================================ S3 배경
s = slide('연구 배경 — 산업 현장의 요통')
headline(s, [('작업 관련 요통은 산업재해의 주요 원인', True, RED, None)],
         top=1.62, size=32)
box, bf = rect(s, 1.3, 2.75, SW - 2.6, 3.35, fill=BOXBG)
bullets(bf, ['물건 들기·운반은 요추 부하의 대표적 원인 작업',
             '고령 근로자·간병 노동 증가로 근골격계 부담 확대',
             '능동 보조 웨어러블이 대안이나, "얼마나 도움 되는가"의 정량 근거가 부족'],
        size=21, sa=26)
bf.vertical_anchor = MSO_ANCHOR.MIDDLE
caption(s, '→ 개발은 활발하지만, 동작별로 얼마나 효과가 있는지에 대한 정량 근거는 상대적으로 부족',
        top=6.35, size=15)
notes(s, '웨어러블 슈트는 많이 개발되지만, 실제로 어느 동작에서 얼마나 '
         '도움이 되는지 정량적으로 제시된 사례는 많지 않습니다.')

# ================================================================ S4 기존 평가 한계
s = slide('기존 평가의 한계')
headline(s, [('표면 근전도(EMG)만으로는 알 수 없는 것들', True, NAVY, None)],
         top=1.58, size=28, h=0.6)
L, W2, T, H2 = 0.75, 5.85, 2.42, 3.55
b1, f1 = rect(s, L, T, W2, H2, fill=BOXBG3, line=RED)
para(f1, 'EMG 측정', size=21, bold=True, color=RED, align=PP_ALIGN.CENTER,
     first=True, space_after=12)
bullets(f1, ['표층 근육만 측정 가능', '개인차·전극 위치에 민감',
             '슈트 착용/미착용 동시 비교 불가', '내부 부하 추정 어려움'],
        size=17, sa=11, indent_marker='–  ')
b2, f2 = rect(s, SW - L - W2, T, W2, H2, fill=BOXBG2, line=GREEN)
para(f2, '근골격계 시뮬레이션', size=21, bold=True, color=GREEN,
     align=PP_ALIGN.CENTER, first=True, space_after=12)
bullets(f2, ['심부 근육 포함 620개 근육 개별 정량', '동일 조건 반복 가능',
             'ON/OFF 동일 동작에서 직접 비교', '관절 모멘트·근육 활성도 직접 산출'],
        size=17, sa=11)
caption(s, '실측과 시뮬레이션은 대체 관계가 아니라 상호 보완 관계', top=6.25, size=15)
notes(s, '실제 실험과 시뮬레이션은 상호 보완적입니다. 시뮬레이션은 EMG로 '
         '측정할 수 없는 심부 근육과, 동일 동작에서의 ON/OFF 비교를 가능하게 합니다.')

# ================================================================ S5 연구 목표
s = slide('연구 목표')
headline(s, [('슈트 효과를, 움직이는 근골격계 모델로 눈에 보이게', True, NAVY, None)],
         top=1.55, size=30, h=0.6)
cards = [('①  정량화', '척추기립근(ES) 활성도를\n슈트 ON / OFF로 직접 비교'),
         ('②  시각화', '근육 부담을 색으로 표현한\n비교 영상 제작'),
         ('③  일반화', '하나의 동작이 아닌,\n5가지 대표 작업 동작으로 확장')]
cw, gap = 3.75, 0.42
x0 = CX - (3 * cw + 2 * gap) / 2
for i, (h1, h2) in enumerate(cards):
    _, cf = rect(s, x0 + i * (cw + gap), 2.48, cw, 2.45, fill=BOXBG, line=NAVY)
    para(cf, h1, size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
         first=True, space_after=14)
    for j, ln_ in enumerate(h2.split('\n')):
        para(cf, ln_, size=16, color=DARK, align=PP_ALIGN.CENTER, space_after=2)
_, bf = rect(s, 1.5, 5.30, SW - 3.0, 0.85, fill=NAVY)
para(bf, '"부하 조건이 다른 5개 동작에서 슈트가 언제, 얼마나 작동하는가"',
     size=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER, first=True, space_after=0)
bf.vertical_anchor = MSO_ANCHOR.MIDDLE
notes(s, '목표는 세 가지입니다. 정량화, 시각화, 그리고 일반화입니다. '
         '특히 세 번째, 하나의 동작이 아니라 부하 조건이 다른 다섯 동작으로 '
         '확장한 것이 본 연구의 차별점입니다.')

# ================================================================ S6 슈트 개요
s = slide('SMA 허리 보조 슈트 개요')
spec = [('액추에이터', 'SMA 직물 근육 (Shape Memory Alloy fabric muscle)'),
        ('구동력', '편측 100 N  (양측 합 200 N)'),
        ('경로', '어깨 → 척추기립근 → 대둔근 → 서혜부'),
        ('모멘트 암', '10 ~ 13 cm'),
        ('보조 토크', '약 20 ~ 26 N·m  (해석 조건 24 N·m)')]
_, hf = tb(s, M, 1.62, 6.4, 0.4)
para(hf, '슈트 사양', size=19, bold=True, color=NAVY, first=True, space_after=0)
ty = 2.15
for k, v in spec:
    _, kf = rect(s, M, ty, 1.75, 0.68, fill=BOXBG)
    para(kf, k, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
         first=True, space_after=0)
    kf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _, vf = tb(s, M + 1.95, ty + 0.10, 4.6, 0.52)
    para(vf, v, size=14.5, color=DARK, first=True, space_after=0)
    vf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ty += 0.78
# 경로 개념도 (도형 작도)
DX = 7.55
_, df = tb(s, DX, 1.62, 5.2, 0.4)
para(df, '슈트 인장 경로 개념도', size=19, bold=True, color=NAVY, first=True, space_after=0)
path_nodes = [('어깨 (상부 고정점)', NAVY), ('척추기립근 라인 (흉·요추)', RED),
              ('대둔근', NAVY), ('서혜부 (하부 고정점)', NAVY)]
ny = 2.15
for i, (nm, col) in enumerate(path_nodes):
    _, nf = rect(s, DX, ny, 3.35, 0.60, fill=BOXBG3 if col == RED else BOXBG,
                 line=col)
    para(nf, nm, size=15, bold=(col == RED), color=col, align=PP_ALIGN.CENTER,
         first=True, space_after=0)
    nf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if i < 3:
        down_arrow(s, DX + 1.52, ny + 0.64, 0.30, 0.30, color=LGRAY)
    ny += 1.05
_, af = rect(s, DX + 3.55, 2.15, 2.05, 3.85, fill=None, line=LGRAY)
para(af, '인장력\n100 N × 2', size=15, bold=True, color=RED,
     align=PP_ALIGN.CENTER, first=True, space_after=10)
para(af, '×', size=13, color=GRAY, align=PP_ALIGN.CENTER, space_after=10)
para(af, '모멘트 암\n0.10 ~ 0.13 m', size=15, bold=True, color=NAVY,
     align=PP_ALIGN.CENTER, space_after=10)
para(af, '=', size=13, color=GRAY, align=PP_ALIGN.CENTER, space_after=10)
para(af, '신전 토크\n20 ~ 26 N·m', size=16, bold=True, color=GREEN,
     align=PP_ALIGN.CENTER, space_after=6)
para(af, '해석 조건: 24 N·m', size=12.5, color=GRAY,
     align=PP_ALIGN.CENTER, space_after=0)
af.vertical_anchor = MSO_ANCHOR.MIDDLE
caption(s, '시뮬레이션에서는 흉추 1번(T1)과 골반 사이에 작용하는 순수 토크 커플 24 N·m으로 모델링 '
           '— 모멘트 암 가정에 의존하지 않는 보수적 표현',
        top=6.15, size=14)
notes(s, '시뮬레이션에서는 슈트를 흉추1번과 골반 사이의 순수 토크 커플 '
         '24 N·m으로 모델링했습니다. 이는 모멘트 암 가정에 의존하지 않는 '
         '보수적 표현입니다.')

# ================================================================ S7 SMA 구동 특성
s = slide('SMA 구동 특성 — 에너지 효율')
headline(s, [('상시 구동(Constant)이 On/Off보다 ', False, NAVY, 28),
             ('13배', True, GREEN, 34), (' 효율적', False, NAVY, 28)],
         top=1.58, size=28, h=0.65)
b1, f1 = rect(s, 0.85, 2.45, 5.6, 2.35, fill=BOXBG3, line=RED)
para(f1, 'On / Off 구동', size=20, bold=True, color=RED,
     align=PP_ALIGN.CENTER, first=True, space_after=12)
para(f1, '가열 – 냉각 반복', size=17, color=DARK, align=PP_ALIGN.CENTER, space_after=8)
para(f1, '→ 매 사이클마다 잠열 재투입 필요', size=17, color=DARK,
     align=PP_ALIGN.CENTER, space_after=10)
para(f1, '소비 에너지  13', size=22, bold=True, color=RED,
     align=PP_ALIGN.CENTER, space_after=0)
b2, f2 = rect(s, SW - 0.85 - 5.6, 2.45, 5.6, 2.35, fill=BOXBG2, line=GREEN)
para(f2, 'Constant 구동', size=20, bold=True, color=GREEN,
     align=PP_ALIGN.CENTER, first=True, space_after=12)
para(f2, '50 ℃ 유지', size=17, color=DARK, align=PP_ALIGN.CENTER, space_after=8)
para(f2, '→ 잠열 재투입 불필요', size=17, color=DARK,
     align=PP_ALIGN.CENTER, space_after=10)
para(f2, '소비 에너지  1', size=22, bold=True, color=GREEN,
     align=PP_ALIGN.CENTER, space_after=0)
_, sf = rect(s, 0.85, 5.05, SW - 1.7, 1.0, fill=BOXBG)
para(sf, '가열 사양   전류 20 A · 저항 1 Ω   |   가열 2 s / 냉각 13 s   |   '
         '유지 2 A @ 50 ℃', size=17, bold=True, color=NAVY,
     align=PP_ALIGN.CENTER, first=True, space_after=0)
sf.vertical_anchor = MSO_ANCHOR.MIDDLE
notes(s, '슈트를 상시 켜두는 것이 오히려 효율적이라는 점은, '
         '실제 착용 운용 시나리오 설계에 중요한 근거입니다.')

# ================================================================ S8 방법론 1 모델
s = slide('방법론 ① — 근골격계 모델')
headline(s, [('전신 620개 근육, 척추기립근 76개를 개별 정량', True, NAVY, None)],
         top=1.52, size=26, h=0.55, align=PP_ALIGN.LEFT)
spec8 = [('모델', 'ThoracolumbarFB v2.0 Full Body (OpenSim 4.x)'),
         ('전신 근육', '620 개'),
         ('척추기립근 (ES)', '76 개  —  Iliocostalis,\nLongissimus thoracis pars lumborum / thoracis'),
         ('흉요추', 'T1 – L5 전 분절 굴곡·신전 자유도'),
         ('해석', 'OpenSim 4.6  /  Static Optimization')]
ty = 2.25
for k, v in spec8:
    _, kf = rect(s, M, ty, 2.0, 0.78, fill=BOXBG)
    para(kf, k, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
         first=True, space_after=0)
    kf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _, vf = tb(s, M + 2.2, ty + 0.07, 4.5, 0.68)
    for j, ln_ in enumerate(v.split('\n')):
        para(vf, ln_, size=15 if j == 0 else 13,
             color=DARK if j == 0 else GRAY, first=(j == 0), space_after=1)
    vf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ty += 0.86
_, mf = rect(s, 7.55, 1.95, 5.18, 4.22, fill=RGBColor(0x11, 0x11, 0x11))
pic(s, f'{MEDIA}/model_fullbody.png', 7.75, 2.08, 1.75, 3.95)
_, lf = tb(s, 9.72, 2.55, 2.75, 3.0)
para(lf, '척추기립근 (ES)', size=17, bold=True, color=RGBColor(0xFF, 0xB3, 0x4D),
     first=True, space_after=10)
para(lf, '주황색 = 본 연구가', size=14, color=RGBColor(0xEE, 0xEE, 0xEE), space_after=2)
para(lf, '정량한 허리 근육군', size=14, color=RGBColor(0xEE, 0xEE, 0xEE), space_after=14)
para(lf, '76 개', size=26, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), space_after=14)
para(lf, '전신 620개 근육 중', size=13.5, color=RGBColor(0xCC, 0xCC, 0xCC), space_after=2)
para(lf, '이 76개를 개별 추적', size=13.5, color=RGBColor(0xCC, 0xCC, 0xCC), space_after=0)
_, c8 = tb(s, 7.55, 6.30, 5.18, 0.4)
para(c8, '전신 모델 측면 렌더 — 근육 색이 진할수록 활성도가 높음', size=13,
     color=GRAY, align=PP_ALIGN.CENTER, first=True, space_after=0)
notes(s, '흉요추 전 분절이 개별 자유도를 갖는 모델이라, 허리 굽힘을 '
         '한 관절이 아니라 여러 분절에 분산해 표현할 수 있습니다.')

# ================================================================ S9 방법론 2 파이프라인
s = slide('방법론 ② — 해석 파이프라인')
steps = [('①\n동작 생성 · 리타겟', '합성 동작 또는 실측 보행\n데이터를 모델 좌표로 변환', NAVY),
         ('②\n동작 육안 검증', '시뮬레이션 전 스냅샷으로\n동작 오류 확인 (필수 관문)', RED),
         ('③\n외력 · 지면반력(GRF)', '박스 하중·지면반력을\n물리적으로 정합하게 부여', NAVY),
         ('④\nStatic Optimization', '동일 동작에서 슈트 토크만\n0 / 24 N·m으로 바꿔 산출', NAVY),
         ('⑤\n근육 부담 시각화', '활성도를 색으로 매핑\n(초록=편함 → 빨강=힘듦)', GREEN)]
bw, gp = 2.16, 0.36
x0 = CX - (5 * bw + 4 * gp) / 2
for i, (hd, ds, col) in enumerate(steps):
    x = x0 + i * (bw + gp)
    _, bfm = rect(s, x, 2.05, bw, 1.62,
                  fill=BOXBG3 if col == RED else (BOXBG2 if col == GREEN else BOXBG),
                  line=col)
    for j, ln_ in enumerate(hd.split('\n')):
        para(bfm, ln_, size=19 if j == 0 else 14, bold=True, color=col,
             align=PP_ALIGN.CENTER, first=(j == 0), space_after=4, line=1.15)
    bfm.vertical_anchor = MSO_ANCHOR.MIDDLE
    _, df2 = tb(s, x, 3.85, bw, 1.05)
    for j, ln_ in enumerate(ds.split('\n')):
        para(df2, ln_, size=12.5, color=DARK, align=PP_ALIGN.CENTER,
             first=(j == 0), space_after=2, line=1.25)
    if i < 4:
        arrow(s, x + bw + 0.045, 2.68, gp - 0.09, 0.36, color=LGRAY)
_, wf = rect(s, 1.4, 5.30, SW - 2.8, 0.95, fill=BOXBG3, line=RED)
rich(wf, [('②단계가 관문 — ', True, RED, 18),
          ('동작이 어색하면 그 뒤 계산은 전부 무의미. 모든 동작은 시뮬레이션 전 '
           '육안 검증 통과 후에만 진행', False, DARK, 17)],
     align=PP_ALIGN.CENTER, first=True, space_after=0)
wf.vertical_anchor = MSO_ANCHOR.MIDDLE
notes(s, '②단계가 핵심입니다. 동작이 어색하면 그 뒤 계산은 전부 무의미하기 때문에, '
         '모든 동작은 시뮬레이션 전에 육안 검증을 통과해야만 진행했습니다.')

# ================================================================ S10 방법론 3 지표
s = slide('방법론 ③ — 평가 지표')
headline(s, [('ES peak를 주 지표로 — EMG 문헌과 정렬', True, NAVY, None)],
         top=1.50, size=26, h=0.55, align=PP_ALIGN.LEFT)
b1, f1 = rect(s, 0.75, 2.15, 5.85, 2.15, fill=BOXBG2, line=GREEN)
para(f1, 'ES peak   (주 지표)', size=20, bold=True, color=GREEN, first=True, space_after=10)
bullets(f1, ['76개 근육 중 최대 활성 근육',
             'EMG 문헌의 40~80 % MVC 범위와 정렬',
             '실제 부하를 반영'], size=16, sa=7)
b2, f2 = rect(s, SW - 0.75 - 5.85, 2.15, 5.85, 2.15, fill=BOXBG, line=GRAY)
para(f2, 'ES mean   (보조 지표)', size=20, bold=True, color=GRAY, first=True, space_after=10)
bullets(f2, ['76개 근육 평균',
             '비활성 근육에 희석되어 실제 부담 과소 표현',
             '지표 무관 강건성 확인용으로 병기'], size=16, sa=7)
_, wf = rect(s, 0.75, 4.58, SW - 1.5, 1.62, fill=BOXBG3, line=RED)
para(wf, '⚠  Reserve 액추에이터 관리', size=19, bold=True, color=RED,
     first=True, space_after=8)
para(wf, 'Reserve가 크면 척추 부하를 근육 대신 흡수해 근육 활성도가 과소평가됨.',
     size=16, color=DARK, space_after=4)
rich(wf, [('본 연구는 spine reserve를 ', False, DARK, 16),
          ('1 ~ 2 N·m 이하', True, RED, 16),
          ('로 제한하여 근육이 부하를 담당하도록 함.', False, DARK, 16)],
     space_after=0)
notes(s, '이 reserve 관리가 결과를 바꾼 실제 사례가 걷기 해석에서 있었고, '
         '뒤에서 다시 설명드리겠습니다.')

# ================================================================ S11 방법론 4 모델 정합성
s = slide('방법론 ④ — 모델 정합성 확보')
headline(s, [('사람이 당연히 하는 동작을 모델이 못 하면, 모델을 고쳤다', True, NAVY, None)],
         top=1.48, size=25, h=0.55, align=PP_ALIGN.LEFT)
hdr = ['발견된 결함', '조치', '정량 영향 (ΔES)']
rows = [('전완(forearm) 정의 오류', '비정상 패치 폐기', '시각화 전용 · 영향 없음'),
        ('손목 3자유도 잠김', '시각화 단계에서 해제', '영향 없음 (손 방향 무관)'),
        ('좌우 어깨대(girdle) 비대칭', '어깨대 전체 거울 대칭화', '영향 없음'),
        ('견갑 protraction 부재', '흉쇄관절 2-DOF 추가', '0.029 %p (무시 가능)'),
        ('좌측 팔 관절축 오류', '축 7개 거울 대칭 수정', '≤ 1.1 %p')]
cw3 = [4.35, 3.95, 3.85]
x0 = M
ty = 2.12
for i, hcol in enumerate(hdr):
    _, hf = rect(s, x0 + sum(cw3[:i]), ty, cw3[i] - 0.06, 0.52, fill=NAVY)
    para(hf, hcol, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         first=True, space_after=0)
    hf.vertical_anchor = MSO_ANCHOR.MIDDLE
ty += 0.56
for r in rows:
    for i, cell in enumerate(r):
        _, cf = rect(s, x0 + sum(cw3[:i]), ty, cw3[i] - 0.06, 0.63,
                     fill=BOXBG if i < 2 else BOXBG2)
        para(cf, cell, size=14, color=DARK if i < 2 else GREEN,
             bold=(i == 2), align=PP_ALIGN.CENTER, first=True, space_after=0)
        cf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ty += 0.67
_, bfm = rect(s, M, ty + 0.10, SW - 2 * M, 0.72, fill=BOXBG2, line=GREEN)
para(bfm, '모든 수정은 ES 정량에 영향이 없거나 무시 가능한 수준임을 재검증(regression)으로 확인',
     size=17, bold=True, color=GREEN, align=PP_ALIGN.CENTER, first=True, space_after=0)
bfm.vertical_anchor = MSO_ANCHOR.MIDDLE
notes(s, '이 표는 결과의 신뢰도를 뒷받침합니다. 모델이 동작을 못 할 때 '
         '자세를 억지로 끼워 맞춘 것이 아니라, 모델 자체의 정의 오류를 '
         '찾아 수정하고, 그 수정이 정량 결과를 바꾸지 않음을 확인했습니다.')

# ================================================================ S12 방법론 5 검증 체계
s = slide('방법론 ⑤ — 검증 체계')
headline(s, [('만든 주체와 검증하는 주체를 분리', True, NAVY, None)],
         top=1.50, size=26, h=0.55, align=PP_ALIGN.LEFT)
b1, f1 = rect(s, 0.75, 2.15, 5.55, 2.55, fill=BOXBG3, line=RED)
para(f1, '문제', size=20, bold=True, color=RED, first=True, space_after=10)
para(f1, '생성자가 자기 결과를 검증하면,\n"수치가 맞으니 그림도 맞다"고 판단하게 됨',
     size=16.5, color=DARK, space_after=8)
para(f1, '→ 자세 오류가 정량 결과와 함께 통과', size=16.5, bold=True,
     color=RED, space_after=0)
arrow(s, 6.45, 3.20, 0.45, 0.45, color=LGRAY)
b2, f2 = rect(s, 7.03, 2.15, 5.55, 2.55, fill=BOXBG2, line=GREEN)
para(f2, '해법', size=20, bold=True, color=GREEN, first=True, space_after=10)
para(f2, '독립 검증자가 수치·의도를 모른 채\n렌더 이미지만 보고 판정',
     size=16.5, color=DARK, space_after=8)
para(f2, '→ 자가 검증 편향 구조적으로 차단', size=16.5, bold=True,
     color=GREEN, space_after=0)
_, cf = tb(s, 0.75, 4.92, SW - 1.5, 0.45)
para(cf, '검증 항목 예시', size=16, bold=True, color=NAVY, first=True, space_after=0)
chk = ['손이 물체를 관통하는가', '발바닥 전체가 접지하는가',
       '양손이 대칭인가', '자세가 자연스러운가']
ccw = (SW - 1.5 - 3 * 0.2) / 4
for i, c in enumerate(chk):
    _, kf = rect(s, 0.75 + i * (ccw + 0.2), 5.40, ccw, 0.6, fill=BOXBG)
    para(kf, c, size=14, color=DARK, align=PP_ALIGN.CENTER, first=True, space_after=0)
    kf.vertical_anchor = MSO_ANCHOR.MIDDLE
caption(s, '검증 통과 전 결과를 채택하지 않음 — 전 동작에 적용',
        top=6.20, size=16, color=NAVY, bold=True)
notes(s, '시뮬레이션에서 수치가 맞아도 동작이 부자연스러우면 결과를 '
         '신뢰할 수 없습니다. 이를 구조적으로 막기 위한 절차입니다.')

# ================================================================ S13 5동작 개요
s = slide('5가지 작업 동작 개요')
headline(s, [('부하 조건이 다른 5개 동작 — 슈트가 언제 작동하는가', True, NAVY, None)],
         top=1.52, size=26, h=0.55, align=PP_ALIGN.LEFT)
mo = [('①', '맨몸 스쿼트', '무릎 굽혀 앉기', '저부하', 'squat'),
      ('②', '맨몸 스툽', '허리 굽혀 숙이기', '중부하', 'stoop'),
      ('③', '박스 들기', '20 kg · 낮은 테이블', '고부하 (들기)', 'box'),
      ('④', '맨몸 걷기', '정상 보행', '초저부하', 'gait'),
      ('⑤', '박스 나르기', '20 kg 안고 보행', '고부하 (지속)', 'carry')]
LOADCOL = {'저부하': GREEN, '중부하': RGBColor(0xB7, 0x7A, 0x0B),
           '고부하 (들기)': RED, '초저부하': GRAY, '고부하 (지속)': RED}
cw, gp = 2.28, 0.33
x0 = CX - (5 * cw + 4 * gp) / 2
for i, (num, nm, ds, ld, key) in enumerate(mo):
    x = x0 + i * (cw + gp)
    _, nf = tb(s, x, 2.12, cw, 0.42)
    para(nf, num, size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
         first=True, space_after=0)
    pic(s, f'{MEDIA}/th_{key}.png', x, 2.58, cw, 1.55)
    _, cf2 = rect(s, x, 4.28, cw, 1.70, fill=BOXBG)
    para(cf2, nm, size=17, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
         first=True, space_after=6)
    para(cf2, ds, size=13, color=DARK, align=PP_ALIGN.CENTER, space_after=10)
    para(cf2, ld, size=14, bold=True, color=LOADCOL[ld],
         align=PP_ALIGN.CENTER, space_after=0)
caption(s, '부하가 다른 5개 동작을 선정 — 슈트가 모든 상황에서 동일하게 작동하는지, '
           '부하에 따라 달라지는지를 확인하기 위함',
        top=6.20, size=14)
notes(s, '부하가 다른 5개 동작을 고른 이유는, 슈트가 모든 상황에서 '
         '똑같이 작동하는지 아니면 부하에 따라 달라지는지를 보기 위함입니다.')

# ============================================================ S14 스쿼트 결과
s = slide('동작 ① 맨몸 스쿼트 — 결과')
headline(s, [('37 % ↓', True, GREEN, 44)], top=1.45, size=44, h=0.78)
_, f1 = rect(s, M, 2.30, 5.55, 3.68, fill=BOXBG)
para(f1, '동작 설명', size=18, bold=True, color=NAVY, first=True, space_after=10)
bullets(f1, ['발바닥 전체 접지 · 무릎만 굽혀 하강',
             '요추 중립 유지 (스툽과의 결정적 차이)',
             '양팔 전방'], size=16, sa=9)
para(f1, '슈트 작동 구간 평균 (주 지표)', size=14, color=DARK, space_before=8, space_after=1)
para(f1, '    ES peak  60.4 % → 37.9 %  =  −37 %', size=15, bold=True,
     color=GREEN, space_after=6)
para(f1, '전주기 정점', size=14, color=DARK, space_after=1)
para(f1, '    ES peak  71.7 % → 46.6 %  =  −35 %', size=15, bold=True,
     color=GREEN, space_after=0)
pic(s, f'{MEDIA}/kf_squat.png', 6.35, 2.30, 6.4, 3.68)
_, cf = rect(s, M, 6.04, SW - 2 * M, 0.80, fill=BOXBG2, line=GREEN)
para(cf, '5동작 완전 통일 조건(동일 모델 + tight reserve) 재해석 값. 이전 47 % 수치는 '
         'reserve·모델 조건이 달랐던 값으로 폐기',
     size=15, color=DARK, align=PP_ALIGN.CENTER, first=True, space_after=0)
cf.vertical_anchor = MSO_ANCHOR.MIDDLE
notes(s, '47%와 37%는 서로 다른 시점의 값입니다. 하나만 쓰면 오해가 생기므로 '
         '두 값을 병기했습니다.')

# ============================================================ S15 스쿼트 영상
s = slide('동작 ① 맨몸 스쿼트 — 비교 영상',
          '슈트 착용 전 / 후 — 허리 근육 부담 비교')
movie(s, 'squat', boxw=9.0, boxh=4.35, top=1.78)
caption(s, VLEGEND, top=6.35, size=14, color=DARK)
notes(s, '영상을 클릭하면 재생됩니다. 좌우 허리 부위 색 차이에 주목해 주십시오. '
         '오른쪽 위 확대창에서 슈트 없음 쪽에만 노란색이 나타납니다.')

# ============================================================ S16 스툽 결과
s = slide('동작 ② 맨몸 스툽 — 결과')
headline(s, [('33 % ↓', True, GREEN, 44)], top=1.45, size=44, h=0.78)
_, f1 = rect(s, M, 2.30, 5.55, 3.68, fill=BOXBG)
para(f1, '동작 설명', size=18, bold=True, color=NAVY, first=True, space_after=10)
bullets(f1, ['고관절 힌지 중심 · 무릎은 편 채 굽힘',
             '요추 굴곡 동반 → ES 부하가 스쿼트보다 큼',
             '산업 현장에서 가장 빈번한 위험 자세'], size=16, sa=9)
para(f1, '슈트 작동 구간 평균 (주 지표)', size=14, color=DARK, space_before=8, space_after=1)
para(f1, '    ES peak  65.3 % → 43.7 %  =  −33 %', size=15, bold=True,
     color=GREEN, space_after=6)
para(f1, '전주기 정점', size=14, color=DARK, space_after=1)
para(f1, '    ES peak  70.4 % → 46.3 %  =  −34 %', size=15, bold=True,
     color=GREEN, space_after=0)
pic(s, f'{MEDIA}/kf_stoop.png', 6.35, 2.30, 6.4, 3.68)
_, cf = rect(s, M, 6.04, SW - 2 * M, 0.80, fill=BOXBG2, line=GREEN)
para(cf, '⭐ 미착용 70.4 %는 선행 실측 EMG(69.8 %MVC)와 0.6 %p 차이 — 모델 신뢰도의 직접 근거',
     size=15, color=DARK, align=PP_ALIGN.CENTER, first=True, space_after=0)
cf.vertical_anchor = MSO_ANCHOR.MIDDLE
notes(s, '여기서 이미 부하가 클수록 상대 감소율이 작아지는 경향이 보입니다. '
         '영상에는 31 %로 표기되어 있는데, 31.8 %를 반올림한 동일한 값입니다.')

# ============================================================ S17 스툽 영상
s = slide('동작 ② 맨몸 스툽 — 비교 영상',
          '슈트 착용 전 / 후 — 허리 근육 부담 비교')
movie(s, 'stoop', boxw=8.6, boxh=4.15, top=1.72)
caption(s, VLEGEND, top=5.98, size=14, color=DARK, h=0.35)
caption(s, '※ 영상 내 표기 "31 %"는 최대 굴곡 시점 −31.8 %를 내림한 값 — 슬라이드의 32 %와 동일한 수치',
        top=6.40, size=12.5, color=GRAY, h=0.35)
notes(s, '스툽은 허리를 직접 굽히는 자세라, 슈트가 작동하는 구간이 뚜렷합니다. '
         '영상에 31 %로 표기된 것은 31.8 %를 내림한 값으로, 슬라이드의 32 %와 같은 수치입니다.')

# ============================================================ S18 박스 과제
s = slide('동작 ③ 박스 들기 — 파지 자세 구현의 과제')
headline(s, [('사람에게 당연한 동작이 모델에서는 어려웠다', True, NAVY, None)],
         top=1.45, size=25, h=0.5, align=PP_ALIGN.LEFT)
_, f1 = rect(s, M, 2.05, 5.9, 3.75, fill=BOXBG2, line=GREEN)
para(f1, '해결한 문제들', size=18, bold=True, color=GREEN, first=True, space_after=12)
for t_ in ['손목 잠김 해제 → 손바닥이 박스를 향하도록',
           '좌우 대칭 파지 → 어깨대 전체 거울 대칭',
           '견갑 protraction 추가 → 팔 도달 거리 확보',
           '자세 배분 정정 → 웅크림(crouch)이 아닌 스툽(stoop)',
           '팔꿈치 자연화 → 옆으로 벌어지지 않고 몸 옆으로']:
    rich(f1, [('✔  ', True, GREEN, 16), (t_, False, DARK, 15.5)], space_after=11)
panels = [('grasp_before_front_pad.png', '이전 — 팔꿈치 벌어짐', RED),
          ('grasp_after_front_pad.png', '완성 앞 — 대칭 파지', GREEN),
          ('grasp_after_side_pad.png', '완성 옆 — 자연 stoop', GREEN)]
pw, pgap = 1.88, 0.17
px0 = 6.75
PH_FIX = 2.55   # 세 패널 높이 통일 → 상·하단 정렬
for i, (fn, lb, col) in enumerate(panels):
    x = px0 + i * (pw + pgap)
    with Image.open(f'{MEDIA}/{fn}') as _im:
        _w = PH_FIX * _im.width / _im.height
    s.shapes.add_picture(f'{MEDIA}/{fn}', I(x + (pw - _w) / 2), I(2.12),
                         I(_w), I(PH_FIX))
    _, pf = tb(s, x - 0.12, 4.86, pw + 0.24, 0.42)
    para(pf, lb, size=12, bold=True, color=col, align=PP_ALIGN.CENTER,
         first=True, space_after=0)
_, gf2 = tb(s, px0, 5.34, 3 * pw + 2 * pgap, 0.42)
para(gf2, '박스 20 kg · 테이블 높이 30 cm 조건에서의 최종 파지 자세',
     size=12, color=GRAY, align=PP_ALIGN.CENTER, first=True, space_after=0)
_, cf = rect(s, M, 5.95, SW - 2 * M, 0.82, fill=BOXBG)
para(cf, '무릎 높이 물체는 쪼그려 앉는 것이 아니라 허리를 굽혀 잡는 것 — '
         '실제 동작 관찰이 시뮬레이션 수렴의 열쇠였음',
     size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER, first=True, space_after=0)
cf.vertical_anchor = MSO_ANCHOR.MIDDLE
notes(s, '이 슬라이드는 학회에서 질문이 나올 만한 부분입니다. '
         '모델이 자동으로 풀지 못한 부분을 어떻게 해결했는지를 보여줍니다.')

# ============================================================ S19 박스 결과
s = slide('동작 ③ 박스 들기 — 결과')
headline(s, [('23 % ↓', True, GREEN, 44),
             ('     ES peak  71.0 %  →  55.0 %', False, DARK, 22)],
         top=1.45, size=44, h=0.78)
_, f1 = rect(s, M, 2.35, 5.55, 3.45, fill=BOXBG)
para(f1, '조건', size=18, bold=True, color=NAVY, first=True, space_after=10)
bullets(f1, ['박스 20 kg, 테이블 높이 30 cm, 박스 30 cm',
             '파지점 약 45 cm, 자연 스툽 자세',
             '슈트 작동 구간(2.2–5.8 s) 평균'], size=16, sa=10)
para(f1, '주 지표 (슈트 작동 구간 평균)', size=14, color=DARK,
     space_before=8, space_after=1)
para(f1, '    ES peak  71.0 % → 55.0 %  =  −22.5 %', size=15, bold=True,
     color=GREEN, space_after=0)
pic(s, f'{MEDIA}/fig_box_es.png', 6.35, 2.35, 6.4, 3.45)
_, cf = rect(s, M, 5.92, SW - 2 * M, 0.85, fill=BOXBG2, line=GREEN)
para(cf, '맨몸 스툽(33 %)보다 감소율이 작음 — 20 kg 하중이 더해져 허리 부하가 커진 만큼, '
         '고정된 24 N·m 보조의 비중이 줄어든 결과',
     size=15, color=DARK, align=PP_ALIGN.CENTER, first=True, space_after=0)
cf.vertical_anchor = MSO_ANCHOR.MIDDLE
notes(s, '절대적인 도움의 크기는 오히려 커지지만, 비율로는 작아집니다.')

# ============================================================ S20 박스 영상
s = slide('동작 ③ 박스 들기 — 비교 영상',
          '20 kg 박스, 테이블 높이 30 cm — 허리 근육 부담 비교')
movie(s, 'box', boxw=8.9, boxh=4.10, top=1.72)
caption(s, VLEGEND, top=5.96, size=14, color=DARK, h=0.35)
caption(s, '※ 영상 내 막대는 반올림 표기(37 % / 29 %) — 정확값은 37.5 % → 28.8 % = −23 %',
        top=6.40, size=12.5, color=GRAY, h=0.35)
notes(s, '박스를 잡고 들어올리는 구간에서 좌우 색 차이가 가장 큽니다.')

# ============================================================ S21 걷기 결과
s = slide('동작 ④ 맨몸 걷기 — 결과')
headline(s, [('감소가 아니라 ', False, DARK, 34), ('재분배', True, RED, 40)],
         top=1.45, size=40, h=0.78)
_, f1 = rect(s, M, 2.35, 5.55, 3.45, fill=BOXBG)
para(f1, '조건', size=18, bold=True, color=NAVY, first=True, space_after=8)
bullets(f1, ['실측 보행 데이터 리타겟 + 실측 지면반력',
             '좌우 팔 교대 스윙 구현'], size=16, sa=8)
para(f1, '총량은 줄고 최대 근육은 늘어남', size=15, bold=True, color=NAVY,
     space_before=6, space_after=6)
for t_, v_, c_ in [('76근육 평균', '−11.9 %', GREEN),
                   ('최대 활성 근육', '+21.4 %', RED),
                   ('장늑근(IL) 24개', '−90.9 %', GREEN),
                   ('최장근 요추부(LTpL)', '+28.5 %', RED)]:
    rich(f1, [('    ' + t_ + '   ', False, DARK, 14), (v_, True, c_, 15)],
         space_after=3)
pic(s, '/data/wearable-assist/opensim_analysis/thoracolumbar_fb/docs/images/'
    'paper_five_motion/fig7_gait_redistribution.png', 6.35, 2.42, 6.4, 3.30)
_, cf = rect(s, M, 5.88, SW - 2 * M, 0.92, fill=BOXBG2, line=GREEN)
rich(cf, [('⭐  부하가 준 것이 아니라 옮겨간 것', True, RED, 16),
          (' — 총량 −17.4 %p이나 최심부 요추근에 집중. 유해 여부는 실측 EMG 검증 필요',
           False, DARK, 15)],
     align=PP_ALIGN.CENTER, first=True, space_after=0)
cf.vertical_anchor = MSO_ANCHOR.MIDDLE
notes(s, '이전 판에서는 무영향이라고 말씀드렸는데, 조건을 통일해 다시 보니 '
         '부하가 줄어든 게 아니라 근육 사이에서 옮겨간 것이었습니다. 총량은 줄지만 '
         '가장 많이 쓰는 근육의 부담은 오히려 늘어납니다. 유해한지는 실측으로 확인해야 합니다.')

# ============================================================ S22 걷기 reserve
s = slide('동작 ④ 걷기 — 해석 지표의 함정',
          'Reserve 설정이 결론을 바꾼 사례')
hdr = ['', '표준 reserve', '정확 reserve (tight)']
rows = [('spine reserve 최대', '16.8 N·m', '1.0 N·m'),
        ('걷기 ES peak (슈트 OFF)', '11 %', '35 %'),
        ('슈트 효과 (전주기 기준)', '−5.6 %p  ("보조")', '−1.0 %p'),
        ('스툽 미착용 ES peak', '31.9 % (문헌 미달)', '70.4 % (문헌 69.8 %)')]
cw3 = [4.5, 3.9, 3.9]
x0 = CX - sum(cw3) / 2
ty = 2.05
for i, hcol in enumerate(hdr):
    _, hf = rect(s, x0 + sum(cw3[:i]), ty, cw3[i] - 0.06, 0.58,
                 fill=NAVY if i == 0 else (RED if i == 1 else GREEN))
    para(hf, hcol, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         first=True, space_after=0)
    hf.vertical_anchor = MSO_ANCHOR.MIDDLE
ty += 0.62
for r in rows:
    for i, cell in enumerate(r):
        _, cf2 = rect(s, x0 + sum(cw3[:i]), ty, cw3[i] - 0.06, 0.72,
                      fill=BOXBG if i == 0 else (BOXBG3 if i == 1 else BOXBG2))
        para(cf2, cell, size=16 if i == 0 else 18,
             bold=(i > 0), color=DARK if i == 0 else (RED if i == 1 else GREEN),
             align=PP_ALIGN.CENTER, first=True, space_after=0)
        cf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    ty += 0.76
_, wf = rect(s, M, ty + 0.18, SW - 2 * M, 1.55, fill=BOXBG3, line=RED)
para(wf, '표준 설정에서는 reserve가 척추 부하를 대신 흡수하여 근육 활성도를 '
         '3배 과소평가하고, 존재하지 않는 "보조 효과"를 만들어냈음.',
     size=16.5, color=DARK, first=True, space_after=6)
para(wf, '※ 위 값은 보행 전주기 기준. 구간별로는 +0.9 ~ +4.3 %p까지 부호가 갈리며(S21), '
         '어느 쪽이든 크기는 들기 동작(23~47 %)에 비해 무시할 수준.',
     size=14.5, color=GRAY, space_after=6)
rich(wf, [('→  저부하 동작일수록 reserve 설정에 민감.', True, RED, 17)], space_after=4)
para(wf, '※ 이후 5동작 전체를 동일 모델·동일 tight 설정으로 재해석하여 조건을 통일함. '
         'tight 조건의 스툽 미착용 70.4 %는 선행 실측 EMG 69.8 %MVC와 0.6 %p 차이로, '
         'tight 설정이 옳았다는 독립 근거.',
     size=13.5, color=GRAY, space_after=0)
notes(s, '방법론적으로 가장 중요한 발견 중 하나입니다. 저부하 조건에서 '
         'reserve를 점검하지 않으면 잘못된 결론에 도달할 수 있습니다.')

# ============================================================ S23 걷기 영상
s = slide('동작 ④ 맨몸 걷기 — 비교 영상',
          '슈트 착용 전 / 후 — 허리 근육 부담 비교')
movie(s, 'gait', boxw=8.9, boxh=4.05, top=1.72)
_, cf = rect(s, 1.6, 5.86, SW - 3.2, 0.80, fill=BOXBG, line=GRAY)
para(cf, '왼쪽 = 슈트 없음   |   오른쪽 = 슈트 착용', size=15, bold=True,
     color=NAVY, align=PP_ALIGN.CENTER, first=True, space_after=3)
para(cf, '좌우 색이 비슷 — 총량은 줄지만 최대 활성 근육은 오히려 증가 (재분배)',
     size=15, color=DARK, align=PP_ALIGN.CENTER, space_after=0)
cf.vertical_anchor = MSO_ANCHOR.MIDDLE
caption(s, '※ 영상 내 22 % / 27 %는 그 순간의 순시값(구간 대표값과 정의가 다름). '
           '자막의 "최대 47 %↓"는 맨몸 스쿼트 기준값',
        top=6.62, size=12.5, color=GRAY, h=0.35)
notes(s, '다른 동작 영상과 달리 좌우 색이 비슷한 것이 정상이며, 이것이 이 동작의 결론입니다. '
         '영상에 뜨는 22 %, 27 %는 그 프레임 순간의 값이고, 앞 슬라이드의 35 %는 '
         '주기 전체의 최대값이라 숫자가 다릅니다.')

# ============================================================ S24 나르기 결과
s = slide('동작 ⑤ 박스 나르기 — 결과')
headline(s, [('약 25 % ↓', True, GREEN, 44)], top=1.45, size=44, h=0.78)
_, f1 = rect(s, M, 2.35, 5.55, 3.30, fill=BOXBG)
para(f1, '조건', size=18, bold=True, color=NAVY, first=True, space_after=8)
bullets(f1, ['20 kg 박스를 배 앞에 안고 보행',
             '전방 하중 보상을 위한 체간 후방 경사 5°'], size=16, sa=8)
para(f1, 'mid-stance ES peak  (대표 지표)', size=14, color=DARK,
     space_before=8, space_after=1)
para(f1, '    99.97 % → 74.54 %  =  −25.4 %p', size=15, bold=True,
     color=GREEN, space_after=6)
para(f1, 'ES mean  (76근육 평균)', size=14, color=DARK, space_after=1)
para(f1, '    18.94 % → 13.76 %  =  −27.4 %', size=15, bold=True,
     color=GREEN, space_after=0)
pic(s, f'{MEDIA}/fig_carry_es.png', 6.35, 2.35, 6.4, 3.30)
_, cf = rect(s, M, 5.78, SW - 2 * M, 1.02, fill=BOXBG3, line=RED)
rich(cf, [('⚠  슈트 미착용 시 최대 활성 근육이 100 %에 포화 — 실제 부담은 그 이상.',
           True, RED, 16)], first=True, space_after=4)
para(cf, '따라서 보고 값은 슈트 효과의 하한선이며, 대표 지표로 mid-stance와 평균값을 사용',
     size=15, color=DARK, space_after=0)
notes(s, '포화 때문에 전체 주기 peak는 효과를 과소평가하므로, '
         'mid-stance와 평균값을 대표 지표로 사용했습니다.')

# ============================================================ S25 나르기 영상
s = slide('동작 ⑤ 박스 나르기 — 비교 영상',
          '20 kg 박스를 안고 보행 — 허리 근육 부담 비교')
movie(s, 'carry', boxw=8.9, boxh=4.05, top=1.72)
_, cf = rect(s, 1.6, 5.92, SW - 3.2, 0.62, fill=BOXBG2, line=GREEN)
para(cf, '그냥 걸을 때는 거의 그대로  —  무거운 것을 안고 걸으면 25 % 감소',
     size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER, first=True, space_after=0)
cf.vertical_anchor = MSO_ANCHOR.MIDDLE
caption(s, '※ 영상 내 막대(예: 85 % → 58 %)는 그 순간의 순시값 — '
           '슬라이드의 25 %는 mid-stance 구간 대표값',
        top=6.62, size=12.5, color=GRAY, h=0.35)
notes(s, '걷기와 나란히 보면 슈트의 선택적 작동 특성이 가장 잘 드러납니다. '
         '영상 막대는 순간값이라 구간을 대표하는 25 %와는 다르게 보일 수 있습니다.')

# ============================================================ S26 통합 결과 (차트)
s = slide('통합 결과 — 부하·효과 패턴')
headline(s, [('슈트는 척추에 부하가 있을 때만 작동한다', True, NAVY, None)],
         top=1.42, size=27, h=0.52, align=PP_ALIGN.LEFT)
cd = CategoryChartData()
cd.categories = ['맨몸 걷기', '박스 들기', '박스 나르기', '맨몸 스툽', '맨몸 스쿼트']
cd.add_series('ES 감소율 (%)', (0, 22.5, 27.6, 33.0, 37.3))
gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, I(0.7), I(2.02),
                        I(SW - 1.4), I(3.62), cd)
ch = gf.chart
ch.has_legend = False
ch.has_title = False
pl = ch.plots[0]
pl.gap_width = 55
pl.has_data_labels = True
dl = pl.data_labels
dl.number_format = '0"%"'
dl.number_format_is_linked = False
dl.position = XL_LABEL_POSITION.OUTSIDE_END
dl.font.size = Pt(18)
dl.font.bold = True
dl.font.name = FONT
dl.font.color.rgb = DARK
BARCOL = [GRAY, RGBColor(0x7F, 0xC2, 0x9B), RGBColor(0x62, 0xB2, 0x86),
          RGBColor(0x35, 0x9A, 0x66), GREEN]
ser = pl.series[0]
for i, c in enumerate(BARCOL):
    pt = ser.points[i]
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = c
    pt.format.line.fill.background()
va = ch.value_axis
va.maximum_scale = 45.0
va.minimum_scale = 0.0
va.has_major_gridlines = True
va.tick_labels.font.size = Pt(13)
va.tick_labels.font.name = FONT
va.tick_labels.font.color.rgb = GRAY
va.has_title = True
va.axis_title.text_frame.text = '척추기립근(ES) 활성도 감소율  (%)'
va.axis_title.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
va.axis_title.text_frame.paragraphs[0].runs[0].font.name = FONT
va.axis_title.text_frame.paragraphs[0].runs[0].font.color.rgb = DARK
ca = ch.category_axis
ca.tick_labels.font.size = Pt(16)
ca.tick_labels.font.bold = True
ca.tick_labels.font.name = FONT
ca.tick_labels.font.color.rgb = DARK
ca.has_major_gridlines = False
_, cf = rect(s, M, 5.70, SW - 2 * M, 1.14, fill=BOXBG2, line=GREEN)
para(cf, '부하가 걸리는 4개 동작에서 부하가 클수록 상대 감소율이 작아짐 '
         '(스쿼트 37 → 스툽 33 → 나르기 28 → 들기 23 %).',
     size=15, color=DARK, align=PP_ALIGN.CENTER, first=True, space_after=3)
para(cf, '※ 걷기는 감소가 아니라 재분배이므로 0으로 표기(S21).  '
         '⚠ 이 단조 경향은 주 지표에서만 성립하며 지표를 바꾸면 성립하지 않음',
     size=12.5, color=GRAY, align=PP_ALIGN.CENTER, space_after=0)
cf.vertical_anchor = MSO_ANCHOR.MIDDLE
notes(s, '이 한 장이 본 연구의 핵심 결과입니다. 위에서 아래로 부하가 커지는 순서인데, '
         '감소율은 반대로 작아집니다. 걷기는 부하 자체가 없어 0에 가깝습니다.')

# ============================================================ S27 고찰 1
s = slide('고찰 ① — 왜 부하가 클수록 비율이 작아지는가')
_, f1 = rect(s, M, 1.72, 6.0, 3.65, fill=BOXBG, line=NAVY)
para(f1, '고정 보조 토크의 구조', size=19, bold=True, color=NAVY,
     first=True, space_after=14)
para(f1, '슈트 보조 토크는  24 N·m 로 고정', size=17, bold=True, color=DARK,
     space_after=12)
para(f1, '↓', size=16, color=LGRAY, space_after=10)
para(f1, '허리 요구 토크가 커질수록\n24 N·m의 상대 비중 감소', size=17,
     color=DARK, space_after=12)
para(f1, '↓', size=16, color=LGRAY, space_after=10)
rich(f1, [('상대 감소율(%) 하락', True, RED, 17),
          (',  그러나 ', False, DARK, 17),
          ('절대 감소량은 유지·증가', True, GREEN, 17)], space_after=0)
_, f2 = rect(s, 7.05, 1.72, SW - M - 7.05, 3.65, fill=BOXBG2, line=GREEN)
para(f2, '예시 계산 (개념)', size=19, bold=True, color=GREEN, first=True, space_after=16)
for req, pct in [('50 N·m', '48 %'), ('100 N·m', '24 %')]:
    rich(f2, [('요구 토크 ' + req, False, DARK, 17),
              ('  →  24 N·m 보조  =  ', False, GRAY, 15),
              (pct + ' 경감', True, NAVY, 19)], space_after=18)
para(f2, '보조량은 같지만, 분모가 2배가 되면\n비율은 절반이 됩니다.',
     size=15, color=GRAY, space_after=0)
_, cf = rect(s, M, 5.55, SW - 2 * M, 1.05, fill=NAVY)
para(cf, '고부하 작업에서 더 큰 비율의 효과를 원한다면 보조 토크 증대가 필요',
     size=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER, first=True, space_after=4)
para(cf, '— 슈트 설계 사양 결정의 직접 근거 —', size=15,
     color=RGBColor(0xCF, 0xDD, 0xE8), align=PP_ALIGN.CENTER, space_after=0)
cf.vertical_anchor = MSO_ANCHOR.MIDDLE
notes(s, '이 관계는 슈트 스펙을 정할 때 바로 쓸 수 있는 설계 지침입니다.')

# ============================================================ S28 고찰 2
s = slide('고찰 ② — 선택적 보조라는 특성')
headline(s, [('필요할 때 돕고, 평상시엔 거스르지 않는다', True, NAVY, None)],
         top=1.55, size=29, h=0.6)
rowsp = [('맨몸 걷기', '거의 무영향', '착용한 채 이동해도 부담 없음', GRAY, BOXBG),
         ('20 kg 나르기', '25 % 감소', '하중이 걸리면 즉시 작동', GREEN, BOXBG2)]
ty = 2.52
for nm, eff, mean, col, bg in rowsp:
    _, f1 = rect(s, 1.05, ty, 3.1, 1.15, fill=bg, line=col)
    para(f1, nm, size=20, bold=True, color=col, align=PP_ALIGN.CENTER,
         first=True, space_after=0)
    f1.vertical_anchor = MSO_ANCHOR.MIDDLE
    arrow(s, 4.32, ty + 0.40, 0.62, 0.36, color=LGRAY)
    _, f2 = rect(s, 5.10, ty, 3.1, 1.15, fill=bg, line=col)
    para(f2, eff, size=22, bold=True, color=col, align=PP_ALIGN.CENTER,
         first=True, space_after=0)
    f2.vertical_anchor = MSO_ANCHOR.MIDDLE
    arrow(s, 8.37, ty + 0.40, 0.62, 0.36, color=LGRAY)
    _, f3 = rect(s, 9.15, ty, SW - M - 9.15, 1.15, fill=bg, line=col)
    para(f3, mean, size=16, color=DARK, align=PP_ALIGN.CENTER,
         first=True, space_after=0)
    f3.vertical_anchor = MSO_ANCHOR.MIDDLE
    ty += 1.42
_, cf = rect(s, 1.05, 5.55, (SW - M) - 1.05, 1.05, fill=BOXBG2, line=GREEN)
para(cf, '상시 구동(Constant) 방식이 에너지 효율적이면서도 일상 동작을 방해하지 않는다는 점이 '
         '실사용 운용 시나리오를 뒷받침',
     size=16.5, bold=True, color=GREEN, align=PP_ALIGN.CENTER, first=True, space_after=0)
cf.vertical_anchor = MSO_ANCHOR.MIDDLE
notes(s, 'S7의 상시 구동 효율성과 연결되는 결론입니다. 상시 켜두어도 '
         '보행을 방해하지 않으므로, 켜고 끄는 제어 없이 운용할 수 있습니다.')

# ============================================================ S29 고찰 3 문헌
# 2026-07-30 전면 재작성: 이전 버전은 선행 연구의 %MVC 절대 포인트를
# 상대 감소율로 오독하여 "본 연구 스쿼트가 선행 범위 초과"라고 서술했음.
s = slide('고찰 ③ — 선행 연구 대조: 본 연구는 보수적 추정')
hdr = ['비교 조건', '본 연구 (ES peak 상대 감소율)', '선행 연구 (실측 EMG)', '대조 결과']
rows = [('맨몸 스툽', '28 ~ 32 % ↓',
         '69.8 → 42.4 %MVC\n= 상대 39.3 % ↓', '본 연구가 더 보수적', GREEN),
        ('맨몸 스쿼트', '37 ~ 47 % ↓',
         '보조 수준 간\n유의차 미보고', '대조 불가', GRAY)]
cw4 = [3.15, 3.35, 3.85, 2.20]
x0 = CX - sum(cw4) / 2
ty = 1.78
for i, hcol in enumerate(hdr):
    _, hf = rect(s, x0 + sum(cw4[:i]), ty, cw4[i] - 0.06, 0.72, fill=NAVY)
    para(hf, hcol, size=14.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         first=True, space_after=0)
    hf.vertical_anchor = MSO_ANCHOR.MIDDLE
ty += 0.76
for r in rows:
    for i, cell in enumerate(r[:4]):
        _, cf2 = rect(s, x0 + sum(cw4[:i]), ty, cw4[i] - 0.06, 0.88,
                      fill=BOXBG if i < 3 else (BOXBG2 if r[4] == GREEN else BOXBG))
        for j, ln_ in enumerate(cell.split('\n')):
            para(cf2, ln_, size=14.5 if i != 3 else 13.5,
                 bold=(i in (1, 3)), color=r[4] if i == 3 else DARK,
                 align=PP_ALIGN.CENTER, first=(j == 0), space_after=1)
        cf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    ty += 0.92
_, nf = rect(s, x0, ty + 0.12, sum(cw4) - 0.06, 1.50, fill=BOXBG2, line=GREEN)
para(nf, '해석', size=16, bold=True, color=GREEN, first=True, space_after=6)
rich(nf, [('직접 대조 가능한 스툽 조건에서 본 연구 감소율이 선행 실측치보다 작음', True, DARK, 14.5),
          (' → 본 시뮬레이션이 슈트 효과를 과대평가하지 않았음을 시사.', False, DARK, 14.5)],
     space_after=4)
rich(nf, [('스쿼트는 대응 선행값이 없어 외부 대조 불가', True, GRAY, 14.5),
          (' — 5동작 중 감소율이 가장 커 실측 EMG 검증의 최우선 대상.', False, DARK, 14.5)],
     space_after=0)
_, rf = tb(s, x0, ty + 1.74, sum(cw4), 0.85)
para(rf, '출처 (원문 초록 직접 확인)', size=12, bold=True, color=GRAY, first=True, space_after=3)
para(rf, 'Hasenmaier et al. (2026) Front Bioeng Biotechnol, doi:10.3389/fbioe.2026.1631785 — n=17, Apogee 능동 외골격',
     size=11.5, color=GRAY, space_after=2)
para(rf, '※ 원문의 "10–27 % MVC"는 %MVC 절대 포인트이며 상대 감소율이 아님 (환산 시 스툽 −39.3 %)',
     size=11.5, color=GRAY, space_after=0)
notes(s, '이 슬라이드는 2026-07-30에 정정했습니다. 이전 판은 선행 연구의 %MVC 절대 '
         '포인트를 상대 감소율로 잘못 읽어 "우리 스쿼트가 문헌보다 크다"고 서술했는데, '
         '원문을 다시 확인하니 스툽 기준으로는 오히려 우리 값이 더 작습니다. '
         '스쿼트는 원문이 유의차를 보고하지 않아 비교 자체가 불가능합니다.')

# ============================================================ S30 고찰 4
s = slide('고찰 ④ — 방법론적 시사점')
pts = [('①', '저부하 동작일수록 reserve 설정에 민감',
        'reserve 점검 없이는 결론이 뒤바뀔 수 있음 (걷기 사례: 효과 −5.6 %p → −1.0 %p)'),
       ('②', '지표 선택이 결과 해석을 좌우',
        'ES peak(EMG 문헌 정렬)를 주 지표로, ES mean은 강건성 확인용 보조 지표로 병기'),
       ('③', '동작의 자연스러움 검증이 정량 신뢰도의 전제',
        '생성 주체와 검증 주체를 분리하는 절차가 필요')]
ty = 1.85
for num, hd, ds in pts:
    _, nf = rect(s, M, ty, 0.85, 1.35, fill=NAVY)
    para(nf, num, size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         first=True, space_after=0)
    nf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _, bf2 = rect(s, M + 1.0, ty, SW - 2 * M - 1.0, 1.35, fill=BOXBG)
    para(bf2, hd, size=19, bold=True, color=NAVY, first=True, space_after=8)
    para(bf2, ds, size=15, color=DARK, space_after=0)
    bf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    ty += 1.55
caption(s, '다른 연구자가 유사 해석을 수행할 때 참고할 수 있는 실무적 교훈',
        top=6.42, size=15, h=0.4)
notes(s, '다른 연구자가 유사 해석을 할 때 참고할 수 있는 실무적 교훈입니다.')

# ============================================================ S31 한계
s = slide('한계')
lim = [('성인 남성 1개 체형 조건',
        '성별·연령·체격 확장 미수행'),
       ('보행·나르기의 지면반력은 타 피험자 실측값 사용',
        '교차 피험자 잔차 존재. 단 슈트 ON/OFF 동일 조건이므로 효과 "차이"는 견고'),
       ('Static Optimization 기반',
        '근육 활성 동역학·수축 속도 미반영'),
       ('나르기 조건에서 최대 활성 근육 포화 (100 %)',
        '보고값은 슈트 효과의 하한선'),
       ('스쿼트 조건은 외부 대조가 불가능',
        '대응 선행 연구가 보조 수준 간 유의차를 보고하지 않음 → 실측 EMG 검증 필요'),
       ('reserve 설정이 동작군에 따라 혼재',
        '보행·운반만 tight — 동작 간 절대 활성도 직접 비교는 제한 (효과 방향은 무관하게 유효)')]
ty = 1.62
for hd, ds in lim:
    _, bf2 = rect(s, M, ty, SW - 2 * M, 0.76, fill=BOXBG3)
    rich(bf2, [(hd, True, DARK, 15.5)], first=True, space_after=3)
    para(bf2, ds, size=13.5, color=GRAY, space_after=0)
    bf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    ty += 0.83
caption(s, '한계를 명확히 하는 것이 결과의 신뢰도를 오히려 높인다',
        top=6.62, size=14, color=NAVY, bold=True, h=0.35)
notes(s, '한계를 명확히 하는 것이 결과의 신뢰도를 오히려 높입니다. '
         '특히 스쿼트 수치가 문헌보다 큰 점은 질문이 나올 수 있는 부분이라 '
         '먼저 밝혀 두었습니다.')

# ============================================================ S32 결론
s = slide('결론')
con = [('①', 'SMA 허리 보조 슈트(24 N·m)는 허리에 부하가 걸리는 동작에서\n'
              '척추기립근 부담을 23 ~ 37 % 감소 '
              '(들기 23 · 나르기 28 · 스툽 33 · 스쿼트 37 %)', GREEN),
       ('②', '감소율은 부하에 반비례 — 부하가 클수록 상대 효과 감소,\n'
              '슈트 설계 사양 결정의 정량 근거 확보', NAVY),
       ('③', '정상 보행에서는 부하 감소가 아니라 재분배 —\n'
              '총량 −17 %p이나 최대 활성 근육 +21 % (유해 여부 검증 필요)', NAVY)]
ty = 1.72
for num, txt, col in con:
    _, nf = rect(s, M, ty, 0.95, 1.28, fill=col)
    para(nf, num, size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         first=True, space_after=0)
    nf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _, bf2 = rect(s, M + 1.1, ty, SW - 2 * M - 1.1, 1.28,
                  fill=BOXBG2 if col == GREEN else BOXBG, line=col)
    for j, ln_ in enumerate(txt.split('\n')):
        para(bf2, ln_, size=18.5, bold=(j == 0), color=DARK,
             first=(j == 0), space_after=3, line=1.25)
    bf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    ty += 1.43
_, cf = rect(s, M + 1.1, 6.02, SW - 2 * M - 1.1, 0.70, fill=NAVY)
para(cf, '5개 작업 동작에 대해 근골격계 모델 기반 정량 평가 및 시각화 완료',
     size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, first=True, space_after=0)
cf.vertical_anchor = MSO_ANCHOR.MIDDLE
notes(s, '세 문장으로 요약하면 이렇습니다. 특히 두 번째와 세 번째가 '
         '설계와 운용에 바로 쓰일 수 있는 결론입니다.')

# ============================================================ S33 향후 계획
s = slide('향후 계획')
road = [('①', '다관절 확장', '허리 + 양측 어깨 + 양측 팔꿈치\n동시 보조 효과 평가'),
        ('②', '인체 조건 확장', '성별·연령별(예: 65세 여성)\n근력 조정 모델 적용'),
        ('③', '동적 해석', 'OpenSim Moco 기반\n근육 활성 동역학 반영')]
cw, gp = 3.75, 0.42
x0 = CX - (3 * cw + 2 * gp) / 2
for i, (num, hd, ds) in enumerate(road):
    x = x0 + i * (cw + gp)
    _, bf2 = rect(s, x, 1.95, cw, 2.65, fill=BOXBG, line=NAVY)
    para(bf2, num, size=26, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
         first=True, space_after=10)
    para(bf2, hd, size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER, space_after=12)
    for j, ln_ in enumerate(ds.split('\n')):
        para(bf2, ln_, size=15, color=DARK, align=PP_ALIGN.CENTER, space_after=2)
    if i < 2:
        arrow(s, x + cw + 0.045, 3.10, gp - 0.09, 0.36, color=LGRAY)
_, f = tb(s, M, 5.10, SW - 2 * M, 0.8)
para(f, '감사합니다', size=32, bold=True, color=NAVY,
     align=PP_ALIGN.CENTER, first=True, space_after=0)
_, f = tb(s, M, 6.05, SW - 2 * M, 0.7)
para(f, '한국기계연구원 (KIMM)   박철훈', size=15, color=DARK,
     align=PP_ALIGN.CENTER, first=True, space_after=3)
para(f, 'parkch@kimm.re.kr', size=15, color=GRAY, align=PP_ALIGN.CENTER, space_after=0)
notes(s, '향후 세 방향으로 확장할 계획입니다. 질문 받겠습니다.')

# ---------------------------------------------------------------- 저장
_sh = autofit_all()
print('autofit 축소 박스:', len(_sh), _sh)
prs.save(OUT)
print('SAVED', OUT, os.path.getsize(OUT) / 1e6, 'MB')
print('slides:', len(prs.slides.__iter__.__self__._sldIdLst))
